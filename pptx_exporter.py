from pptx import Presentation, util
from tqdm import tqdm
import yaml
import re
import os
import csv
import collections

"""
ログからpptxを生成するクラス
クラスと言いつつもはや本体
可読性というか変更しやすいようにクラス化しただけ
"""
class Makepptx:
    def __init__(self, setting_file):
        # ymlからパラメータを読み込み(タイトルやログの場所など)
        self.parameters = Parameters(setting_file)
        # プレゼンテーションを開く(テンプレートを指定することでテンプレートのスライドマスターを使える)
        self.prs = Presentation(self.parameters.template)
        # タイトル作成
        self.make_title()

    """
    リセット
    """
    def reset(self):
        self.prs = Presentation(self.parameters.template)
        self.make_title()

    """
    タイトルスライドを生成する関数
    """
    def make_title(self):
        title_slide = self.prs.slides[0]
        title_slide.shapes.title.text = self.parameters.title

    """
    ログファイルの名前をもとにスライドを生成する関数
    この関数の中身は人によって変わると思うので変更推奨
    """
    def make_slide_from_name(self, log_dir_name, title=""):
        # スライドのレイアウト指定
        slide_layout = self.prs.slide_layouts[5]
        # スライドを追加
        slide = self.prs.slides.add_slide(slide_layout)
        # タイトルを入力
        slide.shapes.title.text = title
        # タイトルのフォントサイズを指定
        slide.shapes.title.text_frame.paragraphs[0].font.size = util.Pt(self.parameters.title_font_size)

        # ログが格納されたディレクトリまでのパス
        log_path = os.path.join(self.parameters.log_path, log_dir_name)

        txt = ""
        """テキストボックスの生成"""
        log_file = self.parameters.log_file
        # ログファイルに入力された文字列の解析
        if re.match(r"[a-zA-Z0-9]*\.csv" ,log_file):
            # csvファイルのパス
            path = os.path.join(log_path, log_file)
            # ログに残したパラメータを読み込み
            param_csv = self.read_csv(path)
            txt = self.csv_to_txt(param_csv)
        elif re.match(r"[a-zA-Z0-9]*\.yml" ,log_file) or re.match(r"[a-zA-Z0-9]*\.yaml" ,log_file):
            # csvファイルのパス
            path = os.path.join(log_path, log_file)
            # ログに残したパラメータを読み込み
            param_yml = self.read_yaml(path)
            txt = ""
            txt = self.yaml_to_txt(txt, param_yml)
        
        
        # テキストボックスの位置(top, left)
        txt_pos = (self.parameters.txt_param["pos_top"], self.parameters.txt_param["pos_left"])
        # テキストボックスの大きさ(width, height)
        txt_size = (self.parameters.txt_param["width"], self.parameters.txt_param["height"])
        # テキストボックスを追加
        textbox = slide.shapes.add_textbox(util.Cm(txt_pos[1]), util.Cm(txt_pos[0]), util.Cm(txt_size[0]), util.Cm(txt_size[1]))
        # 段落の追加
        paragraph = textbox.text_frame.add_paragraph()
        # 書き込み
        paragraph.text = txt
        # フォント指定
        paragraph.font.name = self.parameters.txt_param["font"]
        # フォントサイズ指定
        paragraph.font.size = util.Pt(self.parameters.txt_param["font_size"])

        """
        画像の貼り付け
        重なる場合は貼り付ける順番の考慮が必要
        """
        # 1枚目
        # パスを指定
        img_path = os.path.join(log_path, self.parameters.loss_img_name)
        # 位置(top, left)
        img_pos = (self.parameters.img_param["base_pos_top"], self.parameters.img_param["base_pos_left"])
        # 画像の幅(アスペクト比固定で変えてくれるので幅のみ指定)
        img_width = self.parameters.img_param["width"]
        # スライドに貼り付け
        if self.parameters.img_param["height"] == "auto":
            slide.shapes.add_picture(img_path, img_pos[1], img_pos[0], width = img_width)
        else:
            slide.shapes.add_picture(img_path, img_pos[1], img_pos[0], width = img_width, height=self.parameters.img_param["height"])

        # 2枚目
        img_path = os.path.join(log_path, "accuracy.png")
        img_pos = (self.parameters.img_param["base_pos_top"] + self.parameters.img_param["add_pos_top"],
                    self.parameters.img_param["base_pos_left"] + self.parameters.img_param["add_pos_left"])
        if self.parameters.img_param["height"] == "auto":
            slide.shapes.add_picture(img_path, img_pos[1], img_pos[0], width = img_width)
        else:
            slide.shapes.add_picture(img_path, img_pos[1], img_pos[0], width = img_width, height=self.parameters.img_param["height"])


    """
    全ログディレクトリを走査してスライドを生成する関数
    """
    def make_slides(self, slide_title=""):
        log_dirs = self.get_log_dirs(self.parameters.log_path)
        for log_dir in tqdm(log_dirs):
            self.make_slide_from_name(log_dir, self.parameters.log_file, title=slide_title)
        
    """
    ログディレクトリのリストを返す関数
    """
    def get_log_dirs(self, log_path):
        return os.listdir(log_path)

    """
    保存
    """
    def save(self, path="./log.pptx"):
        self.prs.save(path)

    """
    CSVから読み込んだリストから文字列を生成する関数
    """
    def csv_to_txt(self, param_csv):
        txt = ""
        for i in range(len(param_csv[0])):
            # インデントの初期設定
            indent = "\t"

            # 桁数の調整
            # Lossが十分小さくなる場合以下式では対応出来ないため考え中
            #if not param_csv[1][i] == None and not re.findall('[0-9]+[.]{0,1}[0-9]{3}', param_csv[1][i])[0] == "0.000":
            #    param_csv[1][i] = re.findall('[0-9]+[.]{0,1}[0-9]{3}', param_csv[1][i])[0]

            # テキストボックスに書き込む内容を生成
            txt = txt + param_csv[0][i] + indent + "= " + param_csv[1][i] + "\n"
        return txt

    """
    CSVファイルを読み込んでリストで返す関数
    辞書型が面倒だったのでリストにしただけ
    """
    def read_csv(self, csv_path):
        with open(csv_path, encoding='utf-8-sig') as f:
            reader = csv.reader(f)
            l = [row for row in reader]
            return l

    """
    yamlから読み込んだリストから文字列を生成する関数(多次元dict対応)
    """
    def yaml_to_txt(self, txt, param_yml, h=0):
        keys = list(dict.keys())
        for key in keys:
            indent = ""
            for i in range(h):
                indent = indent + "\t"
            if isinstance(param_yml[key], collections.abc.Mapping) or type(param_yml[key]) is dict:
                txt = txt + indent + key + "\n"
                txt = self.yaml_to_txt(txt, param_yml[key], h+1)
            else:
                txt = txt + indent + key + str(param_yml[key]) + "\n"
                # 桁数の調整
                # Lossが十分小さくなる場合以下式では対応出来ないため考え中
                #if not param_yml[key] == None and not re.findall('[0-9]+[.]{0,1}[0-9]{3}', param_yml[key])[0] == "0.000":
                #    txt = txt + indent + key + "{:.3f}".format(str(param_yml[key])) + "\n"
        return txt

    def read_yaml(setting_yaml_file):
        with open(setting_yaml_file) as f:
            return yaml.safe_load(f)

"""
ymlからパラメータを読み込んで保持しておくクラス
"""
class Parameters:
    def __init__(self, yaml_file):
        all_parameters = self.read_yaml(yaml_file)
        self.title = all_parameters["name"]
        self.log_path = all_parameters["log_path"]
        self.template = all_parameters["template"]
        self.txt_param = all_parameters["txt_box"]
        self.img_param = all_parameters["img"]
        self.acc_img_name = all_parameters["acc_img_name"]
        self.loss_img_name = all_parameters["loss_img_name"]
        self.log_file = all_parameters["log_file"]
        self.title_font_size = all_parameters["title_font_size"]
    
    def read_yaml(self, yaml_file):
        with open(yaml_file) as f:
            return yaml.safe_load(f)


"""
実行用関数
"""
def main():
    wrapper = Makepptx()
    wrapper.save()


if __name__ == '__main__':
    main()